#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate anteriq-pitch.pptx from the current HTML content. All text is verbatim."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

INK = RGBColor(0x1A, 0x1A, 0x2E)
INK_SOFT = RGBColor(0x3D, 0x3D, 0x54)
INK_MUTED = RGBColor(0x7A, 0x7A, 0x8C)
PAPER = RGBColor(0xF8, 0xF6, 0xF1)
PAPER_WARM = RGBColor(0xF0, 0xEC, 0xE3)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
BRONZE = RGBColor(0x7A, 0x68, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_DARK = RGBColor(0x25, 0x25, 0x3E)
TXT_L = RGBColor(0xA0, 0x9E, 0x96)
HL_BG = RGBColor(0xE8, 0xE0, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.7)
CW = SLIDE_W - 2 * M

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BL = prs.slide_layouts[6]

def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c
def txb(s, l, t, w, h):
    return s.shapes.add_textbox(l, t, w, h)
def sr(tf, txt, sz=14, b=False, it=False, c=INK_SOFT, al=PP_ALIGN.LEFT):
    tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = al
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.bold = b
    r.font.italic = it; r.font.color.rgb = c; return r
def ap(tf, txt, sz=14, b=False, it=False, c=INK_SOFT, al=PP_ALIGN.LEFT, sb=Pt(4)):
    p = tf.add_paragraph(); p.alignment = al; p.space_before = sb
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.bold = b
    r.font.italic = it; r.font.color.rgb = c; return p

def hdr(s, act, lab, title, sub, dark=False):
    bg(s, INK if dark else PAPER)
    tc = PAPER if dark else INK; sc = TXT_L if dark else INK_SOFT; lc = GOLD if dark else BRONZE
    y = Inches(0.4)
    if act:
        t = txb(s, M, y, CW, Inches(0.3)); sr(t.text_frame, act, 9, True, c=lc); y += Inches(0.28)
    if lab:
        t = txb(s, M, y, CW, Inches(0.3)); sr(t.text_frame, lab.upper(), 11, True, c=lc); y += Inches(0.35)
    t = txb(s, M, y, CW, Inches(0.6)); sr(t.text_frame, title, 28, True, c=tc); y += Inches(0.65)
    if sub:
        t = txb(s, M, y, CW, Inches(0.6)); sr(t.text_frame, sub, 16, c=sc); y += Inches(0.65)
    return y

def cbox(s, l, t, w, h, bg_c=PAPER_WARM, bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_c
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(2)
    else: sh.line.fill.background()
    sh.shadow.inherit = False; return sh

def ctxt(s, l, t, w, h, tit, bod, extra=None, bg_c=PAPER_WARM, bc=None, tc=INK, bc2=INK_SOFT):
    cbox(s, l, t, w, h, bg_c, bc); p = Inches(0.15)
    tb = txb(s, l+p, t+p, w-2*p, h-2*p); tf = tb.text_frame; tf.word_wrap = True
    sr(tf, tit, 13, True, c=tc); ap(tf, bod, 11, c=bc2)
    if extra: ap(tf, extra, 10, True, c=tc)

def punch(s, y, txt, dark=False):
    t = txb(s, M, y, CW, Inches(0.6))
    sr(t.text_frame, txt, 14, it=True, c=TXT_L if dark else INK_SOFT, al=PP_ALIGN.CENTER)


# ===== S1 COVER =====
s = prs.slides.add_slide(BL); bg(s, INK)
t = txb(s, Inches(1.5), Inches(2.0), Inches(10), Inches(0.5))
sr(t.text_frame, "ANTERIQ", 20, True, c=GOLD, al=PP_ALIGN.CENTER)
t = txb(s, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0))
sr(t.text_frame, "La transformation IA\nde bout en bout", 40, True, c=PAPER, al=PP_ALIGN.CENTER)
t = txb(s, Inches(2.5), Inches(4.2), Inches(8), Inches(0.5))
sr(t.text_frame, "Strat\u00e9gie \u00b7 Donn\u00e9es \u00b7 IA \u00b7 Infrastructure souveraine \u00b7 Humain", 16, it=True, c=GOLD, al=PP_ALIGN.CENTER)
t = txb(s, Inches(2), Inches(5.0), Inches(9), Inches(0.7))
sr(t.text_frame, "Le premier groupe capable d\u2019accompagner une organisation du diagnostic strat\u00e9gique au d\u00e9ploiement en production \u2014 et de laisser le client propri\u00e9taire et autonome.", 14, c=TXT_L, al=PP_ALIGN.CENTER)
t = txb(s, Inches(4), Inches(6.2), Inches(5), Inches(0.4))
sr(t.text_frame, "Confidentiel \u2014 Avril 2026", 12, c=INK_MUTED, al=PP_ALIGN.CENTER)

# ===== S2 LE PROBL\u00c8ME =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte I", "Le probl\u00e8me", "Trois options sur le march\u00e9. Toutes insatisfaisantes.", "Les organisations veulent se transformer. Le march\u00e9 ne leur offre que des fragments.", True)
cw = (CW - Inches(0.4)) / 3
for i, (t_, b_) in enumerate([
    ("Les ESN et int\u00e9grateurs", "Ils vendent du temps-homme et des POC qui ne passent jamais en production. Le client reste d\u00e9pendant. Les donn\u00e9es partent dans le cloud. Personne ne s\u2019occupe des \u00e9quipes."),
    ("Les \u00e9diteurs SaaS", "Lock-in maximal, donn\u00e9es hors de contr\u00f4le, personnalisation limit\u00e9e. Le client loue une solution qu\u2019il ne comprend pas et ne poss\u00e8de pas."),
    ("Les cabinets de conseil", "Ils livrent des slides et partent. Pas de d\u00e9ploiement, pas d\u2019infra, pas de suivi. Le plan reste dans un tiroir."),
]):
    ctxt(s, M + i*(cw+Inches(0.2)), y, cw, Inches(1.5), t_, b_, bg_c=CARD_DARK, tc=PAPER, bc2=TXT_L)
sy = y + Inches(1.7)
for i, (n, l, src) in enumerate([
    ("80%+", "des projets IA \u00e9chouent", "RAND Corporation, 2024"),
    ("85%", "des mod\u00e8les IA n\u2019atteignent pas la production", "Gartner"),
    ("8%", "des grandes entreprises ont d\u00e9ploy\u00e9 l\u2019IA \u00e0 l\u2019\u00e9chelle", "McKinsey / Institut de l\u2019Entreprise, 2025"),
]):
    t = txb(s, M + i*(cw+Inches(0.2)), sy, cw, Inches(1.2)); tf = t.text_frame; tf.word_wrap = True
    sr(tf, n, 36, True, c=GOLD, al=PP_ALIGN.CENTER)
    ap(tf, l, 11, c=TXT_L, al=PP_ALIGN.CENTER); ap(tf, src, 8, c=INK_MUTED, al=PP_ALIGN.CENTER)

# ===== S3 NOTRE R\u00c9PONSE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte I", "Notre r\u00e9ponse", "Un seul interlocuteur. La transformation compl\u00e8te.",
        "Du diagnostic strat\u00e9gique au d\u00e9ploiement en production, en passant par la structuration des donn\u00e9es et l\u2019accompagnement des \u00e9quipes.", True)
fx = M + Inches(0.5)
for item in [("D\u00e9couverte","Mod\u00e8le cible de l\u2019entreprise"), None, ("Impl\u00e9mentation","Donn\u00e9es & outils du quotidien"), None, ("Impact humain","Formation, reskilling, transitions")]:
    if item is None:
        t = txb(s, fx, y+Inches(0.15), Inches(0.5), Inches(0.4)); sr(t.text_frame, "\u2192", 24, c=INK_MUTED, al=PP_ALIGN.CENTER); fx += Inches(0.6)
    else:
        bw = Inches(2.5); cbox(s, fx, y, bw, Inches(0.7), CARD_DARK)
        t = txb(s, fx, y+Inches(0.05), bw, Inches(0.6)); tf = t.text_frame; tf.word_wrap = True
        sr(tf, item[0], 13, True, c=PAPER, al=PP_ALIGN.CENTER); ap(tf, item[1], 10, c=TXT_L, al=PP_ALIGN.CENTER); fx += bw + Inches(0.1)
y += Inches(1.0); cw2 = (CW - Inches(0.2)) / 2
for i, (t_, b_) in enumerate([
    ("Souveraine", "Les donn\u00e9es restent chez le client. Pas de d\u00e9pendance cloud, pas de juridiction \u00e9trang\u00e8re. EU AI Act, NIS2, r\u00e9glementations sectorielles."),
    ("\u00c9mancipatrice", "Le client est propri\u00e9taire de tout \u2014 donn\u00e9es, mod\u00e8les, apps, code. La plateforme reste chez lui pour les mises \u00e0 jour et le support."),
    ("Scientifiquement rigoureuse", "Nos propres mod\u00e8les IA, open source, \u00e9valu\u00e9s par les pairs, \u00e9prouv\u00e9s en production. 10 ans de R&D. Pas de bo\u00eetes noires."),
    ("Humaine", "Chaque d\u00e9ploiement inclut un plan humain. Formation, reskilling, coaching. Personne n\u2019est laiss\u00e9 de c\u00f4t\u00e9."),
]):
    col, row = i%2, i//2; left = M + col*(cw2+Inches(0.2)); top = y + row*Inches(1.1)
    ctxt(s, left, top, cw2, Inches(0.95), t_, b_, bg_c=CARD_DARK, tc=GOLD, bc2=TXT_L)

# ===== S4 DISCOVERY =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte II", "Phase 1 \u2014 Discovery", "Construire la tour de contr\u00f4le de la transformation.",
        "L\u2019engagement commence au niveau du comit\u00e9 ex\u00e9cutif. Le plan qu\u2019on construit est directement outillable par la plateforme du groupe.")
cw2 = (CW - Inches(0.3)) / 2
cbox(s, M, y, cw2, Inches(3.5), PAPER_WARM)
t = txb(s, M+Inches(0.15), y+Inches(0.1), cw2-Inches(0.3), Inches(3.3)); tf = t.text_frame; tf.word_wrap = True
sr(tf, "Ce que nous faisons", 13, True, c=INK)
for b in [
    "Acculturation des dirigeants \u2014 lecture strat\u00e9gique de ce que l\u2019IA change dans leur mod\u00e8le d\u2019affaires, leur cha\u00eene de valeur, leur structure de co\u00fbts",
    "Cartographie des opportunit\u00e9s \u2014 cas d\u2019usage par fonction, performance attendue, impact business et humain quantifi\u00e9",
    "Feuille de route budg\u00e9t\u00e9e \u2014 plan prioris\u00e9 avec business cases, s\u00e9quencement r\u00e9aliste, analyse make-or-buy",
    "Fondations de l\u2019ontologie \u2014 premier mod\u00e8le de donn\u00e9es structur\u00e9 \u00e0 partir des processus et flux du client",
]: ap(tf, f"\u2022 {b}", 11, c=INK_SOFT)

rx = M + cw2 + Inches(0.3)
cbox(s, rx, y, cw2, Inches(3.5), PAPER_WARM, BRONZE)
t = txb(s, rx+Inches(0.15), y+Inches(0.1), cw2-Inches(0.3), Inches(3.3)); tf = t.text_frame; tf.word_wrap = True
sr(tf, "Ce qui rend cette phase diff\u00e9rente", 13, True, c=INK)
ap(tf, "Elle n\u2019est pas d\u00e9connect\u00e9e de l\u2019ex\u00e9cution. Il n\u2019y a pas de \u00ab passation \u00bb \u00e0 un int\u00e9grateur tiers \u2014 les m\u00eames personnes qui diagnostiquent sont celles qui accompagnent le d\u00e9ploiement.", 11, c=INK_SOFT)
ap(tf, "", 6); ap(tf, "Briques activ\u00e9es", 11, True, c=BRONZE)
ap(tf, "Conseil strat\u00e9gique Data & IA \u00b7 Ontologie & donn\u00e9es", 11, c=INK_SOFT)
ap(tf, "", 6); ap(tf, "Revenu", 11, True, c=BRONZE)
ap(tf, "One-shot \u00b7 100\u2013300K par programme", 11, c=INK_SOFT)

# ===== S5 BUILD =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte II", "Phase 2 \u2014 Build", "Mettre en place le plan de cr\u00e9ation de valeur. Outiller l\u2019organisation.",
        "Le conseil et le build ne sont pas s\u00e9quentiels \u2014 ils sont int\u00e9gr\u00e9s. R\u00e9sultats tangibles en semaines.")
cw4 = (CW - Inches(0.6)) / 4
for i, (t_, b_) in enumerate([
    ("Infrastructure", "Co-construction avec le DSI. Plateforme Hypsis install\u00e9e on-prem. Pipelines de donn\u00e9es connect\u00e9s. Ontologie structur\u00e9e."),
    ("Applications m\u00e9tier", "Chaque cas d\u2019usage devient un outil en production. Pas des POC \u2014 des outils utilis\u00e9s au quotidien, aliment\u00e9s par l\u2019ontologie et les mod\u00e8les IA."),
    ("Acculturation DSI", "Le DSI est un alli\u00e9. On construit avec lui, on lui transf\u00e8re la comp\u00e9tence. Il devient le gardien de l\u2019\u00e9cosyst\u00e8me IA."),
    ("Formation & reskilling", "Plan humain en parall\u00e8le d\u00e8s le premier jour. Formation, reconversion, coaching des managers, conduite du changement."),
]):
    bc = BRONZE if i == 3 else None
    ctxt(s, M+i*(cw4+Inches(0.2)), y, cw4, Inches(2.0), t_, b_, bg_c=PAPER_WARM, bc=bc)
punch(s, y+Inches(2.3), "Toutes les briques du groupe sont activ\u00e9es. Conseil \u00b7 Ontologie \u00b7 Infrastructure \u00b7 IA \u00b7 Humain. Revenu : one-shot \u00b7 100\u2013500K.")

# ===== S6 RUN =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte II", "Phase 3 \u2014 Run", "Le client est autonome. Le produit reste chez lui.",
        "Le client peut partir avec tout ce qui a \u00e9t\u00e9 construit. C\u2019est pr\u00e9cis\u00e9ment parce qu\u2019il peut partir qu\u2019il choisit de rester.", True)
cw3 = (CW - Inches(0.4)) / 3
for i, (t_, b_) in enumerate([
    ("Plateforme en place", "Hypsis reste d\u00e9ploy\u00e9e chez le client. Les \u00e9quipes internes construisent de nouvelles apps, adaptent les existantes \u2014 sans d\u00e9pendre de nous."),
    ("Intelligence p\u00e9renne", "Mod\u00e8les IA, algorithmes d\u2019optimisation, briques documentaires restent embarqu\u00e9s. Chaque mise \u00e0 jour am\u00e8ne de nouvelles capacit\u00e9s \u2014 sans projet de r\u00e9int\u00e9gration."),
    ("Propri\u00e9t\u00e9 totale", "Donn\u00e9es, mod\u00e8les, applications, code \u2014 tout appartient au client. Rien n\u2019est pi\u00e9g\u00e9 dans un cloud ou derri\u00e8re une API propri\u00e9taire."),
]):
    bc = GOLD if i == 2 else None
    ctxt(s, M+i*(cw3+Inches(0.2)), y, cw3, Inches(1.8), t_, b_, bg_c=CARD_DARK, bc=bc, tc=PAPER, bc2=TXT_L)
punch(s, y+Inches(2.1), "Ce qui g\u00e9n\u00e8re du revenu r\u00e9current : l\u2019abonnement Hypsis \u2014 mises \u00e0 jour continues de la plateforme et des mod\u00e8les dans un environnement qui \u00e9volue vite, formation et certification permanentes des \u00e9quipes client. 3\u201315K/mois par client.", True)

# ===== S7 VERTICALES + BRIQUES =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte III", "Le groupe", "Deux verticales. Cinq briques.",
        "Le groupe s\u2019organise en verticales sectorielles port\u00e9es par des experts m\u00e9tier, rendues possibles par cinq capacit\u00e9s transverses.")
cw2 = (CW - Inches(0.3)) / 2
for i, (t_, b_) in enumerate([
    ("Transformation ops industrielle", "A\u00e9ro, d\u00e9fense, industrie. Partenaire expert du secteur, int\u00e9gr\u00e9 au groupe."),
    ("Transformation ops tertiaire", "Services, finance, retail, luxe. Unit\u00e9 d\u00e9di\u00e9e, constitu\u00e9e par build-up."),
]):
    ctxt(s, M+i*(cw2+Inches(0.3)), y, cw2, Inches(0.9), t_, b_, bg_c=PAPER_WARM, bc=BRONZE)
y += Inches(1.1); cw5 = (CW - Inches(0.5)) / 5
for i, (t_, b_, phase, sub) in enumerate([
    ("Strat\u00e9gie Data & IA", "COMEX, Master Plan, CDO Office, d\u00e9ploiement", "Discovery + Build", ""),
    ("Ontologie & donn\u00e9es", "Cartographie, ing\u00e9nierie des donn\u00e9es, int\u00e9gration legacy", "Discovery + Build", "AI&D \u00b7 30 consultants \u00b7 6M CA"),
    ("Plateforme souveraine", "Cr\u00e9ation, d\u00e9ploiement, supervision d\u2019apps IA on-prem. Capitalisation de l\u2019exp\u00e9rience.", "Build + Run", "Hypsis"),
    ("Excellence scientifique", "Mod\u00e8les IA open source, \u00e9valu\u00e9s par les pairs", "Build + Run", "Jolibrain \u00b7 5 PhDs \u00b7 50+ GPUs"),
    ("Transformation humaine", "Formation, reskilling, conduite du changement", "Build", "Horizon d\u00e9but 2027"),
]):
    left = M + i*(cw5+Inches(0.125)); cbox(s, left, y, cw5, Inches(2.2), PAPER_WARM)
    t = txb(s, left+Inches(0.1), y+Inches(0.08), cw5-Inches(0.2), Inches(2.0)); tf = t.text_frame; tf.word_wrap = True
    sr(tf, t_, 11, True, c=INK); ap(tf, b_, 9, c=INK_SOFT)
    if sub: ap(tf, sub, 8, c=INK_MUTED)
    ap(tf, phase, 9, True, c=INK)
punch(s, y+Inches(2.4), "Les verticales apportent l\u2019expertise m\u00e9tier. Les briques apportent les capacit\u00e9s. C\u2019est l\u2019int\u00e9gration des deux qui produit la transformation.")

# ===== S8 CHA\u00ceNON MANQUANT =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte III", "Le cha\u00eenon manquant", "La brique strat\u00e9gique qui transforme le groupe.",
        "La phase de Discovery est le point d\u2019entr\u00e9e de toute la cha\u00eene de valeur. C\u2019est elle qui convainc le COMEX et g\u00e9n\u00e8re le pipeline.", True)
cw2 = (CW - Inches(0.3)) / 2
cbox(s, M, y, cw2, Inches(3.2), CARD_DARK, GOLD)
t = txb(s, M+Inches(0.15), y+Inches(0.1), cw2-Inches(0.3), Inches(3.0)); tf = t.text_frame; tf.word_wrap = True
sr(tf, "Le profil id\u00e9al", 13, True, c=PAPER)
for b in [
    "Cabinet fran\u00e7ais, ind\u00e9pendant, enti\u00e8rement d\u00e9di\u00e9 \u00e0 la strat\u00e9gie Data & IA",
    "Positionn\u00e9 en amont \u2014 aupr\u00e8s des DG, CDO, responsables m\u00e9tier",
    "Cycle de vie complet formalis\u00e9 : de l\u2019acculturation \u00e0 l\u2019autonomie des BU",
    "Culture de rigueur strat\u00e9gique (70% advisory / 30% build)",
    "30+ grands comptes, taux de renouvellement \u00e9lev\u00e9",
    "\u00c9cosyst\u00e8me \u00e0 fa\u00e7on : universit\u00e9 interne, veille tech, capitalisation des connaissances",
]: ap(tf, f"\u2022 {b}", 10, c=TXT_L)

rx = M + cw2 + Inches(0.3)
cbox(s, rx, y, cw2, Inches(3.2), CARD_DARK)
t = txb(s, rx+Inches(0.15), y+Inches(0.1), cw2-Inches(0.3), Inches(3.0)); tf = t.text_frame; tf.word_wrap = True
sr(tf, "La synergie", 13, True, c=PAPER)
ap(tf, "\u00ab Demain, le conseil strat\u00e9gique ne pourra plus se contenter de recommander \u2014 il devra d\u00e9ployer. Pour d\u00e9ployer, il faut une plateforme, des mod\u00e8les, et une infrastructure. \u00bb", 11, it=True, c=TXT_L)
ap(tf, "", 6); ap(tf, "Ce que le groupe apporte", 11, True, c=GOLD)
ap(tf, "Mod\u00e8les IA propri\u00e9taires, plateforme souveraine, infrastructure \u2014 ce que les cabinets conseil IA sous-traitent aujourd\u2019hui \u00e0 des ESN.", 10, c=TXT_L)
ap(tf, "", 6); ap(tf, "Ce qu\u2019un tel cabinet apporte", 11, True, c=GOLD)
ap(tf, "Acc\u00e8s direct aux COMEX des plus grandes entreprises fran\u00e7aises, 50+ consultants certifi\u00e9s, moteur commercial \u00e0 80% de r\u00e9currence.", 10, c=TXT_L)

# ===== S9 STRUCTURE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte IV", "Structure", "Un groupe con\u00e7u pour le build-up.",
        "Filiales autonomes, P&L s\u00e9par\u00e9s, gouvernance compatible fonds. Chaque fondateur qui rejoint garde les r\u00eanes de ce qu\u2019il a construit.")
sw = (CW - Inches(0.8)) / 5
for i, (n, l) in enumerate([
    ("17M\u20ac", "CA consolid\u00e9 ann\u00e9e 1"), ("95+", "collaborateurs"),
    ("380+", "GPUs"), ("4:1", "ratio GPU / collaborateur"), ("5", "briques de capacit\u00e9"),
]):
    t = txb(s, M+i*(sw+Inches(0.2)), y, sw, Inches(0.9)); tf = t.text_frame
    sr(tf, n, 32, True, c=BRONZE, al=PP_ALIGN.CENTER); ap(tf, l, 11, c=INK_SOFT, al=PP_ALIGN.CENTER)
y += Inches(1.1); cw3 = (CW - Inches(0.4)) / 3
for i, (t_, b_) in enumerate([
    ("Gouvernance", "CEO groupe : Yvan Chabanne, ex-DG Scalian (~5 000 pers.)\nChaque BU a son propre CEO et son P&L\nCEOs entrants obtiennent un si\u00e8ge au conseil"),
    ("Autonomie des BU", "Le groupe ajoute les synergies commerciales et les services partag\u00e9s. Il ne micro-manage pas. Un fondateur qui rejoint garde son CEO, son \u00e9quipe, ses clients."),
    ("BSPCE", "La Loi de Finances 2026 permet \u00e0 la holding d\u2019\u00e9mettre des BSPCE aux employ\u00e9s de toute filiale d\u00e9tenue \u00e0 75%+. Equity dans la holding, avantage fiscal, vesting 4 ans."),
]):
    bc = BRONZE if i == 2 else None
    ctxt(s, M+i*(cw3+Inches(0.2)), y, cw3, Inches(2.0), t_, b_, bg_c=PAPER_WARM, bc=bc)

# ===== S10 MOD\u00c8LE \u00c9CONOMIQUE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte IV", "Mod\u00e8le \u00e9conomique", "Nous vendons de l\u2019intelligence. Humaine ou artificielle.",
        "Dans les deux cas, une partie de cette intelligence reste chez le client \u2014 dans les comp\u00e9tences transf\u00e9r\u00e9es, les mod\u00e8les d\u00e9ploy\u00e9s, les applications qu\u2019il poss\u00e8de. C\u2019est ce d\u00e9p\u00f4t permanent qui cr\u00e9e la valeur et la relation de long terme.", True)
cw3 = (CW - Inches(0.4)) / 3
ctxt(s, M, y, cw3, Inches(1.6), "Intelligence humaine",
     "Jours de conseil, d\u00e9ploiement, formation. One-shot \u00b7 100\u2013300K par programme.",
     bg_c=CARD_DARK, tc=PAPER, bc2=TXT_L)
ctxt(s, M+cw3+Inches(0.2), y, cw3, Inches(1.6), "Intelligence artificielle",
     "Compute, mod\u00e8les IA, mises \u00e0 jour, certification des \u00e9quipes client. R\u00e9current \u00b7 3\u201315K/mois.",
     bg_c=CARD_DARK, tc=PAPER, bc2=TXT_L)
rx = M + 2*(cw3+Inches(0.2))
cbox(s, rx, y, cw3, Inches(1.6), CARD_DARK, GOLD)
t = txb(s, rx+Inches(0.15), y+Inches(0.1), cw3-Inches(0.3), Inches(1.4)); tf = t.text_frame; tf.word_wrap = True
sr(tf, "Trajectoire du mix", 13, True, c=PAPER)
ap(tf, "Aujourd\u2019hui : ~95% humaine / 5% artificielle", 11, c=TXT_L)
ap(tf, "Ann\u00e9e 3 : ~73% humaine / 27% artificielle", 11, c=TXT_L)
ap(tf, "Ann\u00e9e 5 : ~63% humaine / 37% artificielle", 11, c=TXT_L)
ap(tf, "L\u2019intelligence artificielle est le multiplicateur de valorisation.", 10, it=True, c=TXT_L)

# ===== S11 TRAJECTOIRE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte IV", "Trajectoire", "17M \u2192 80M en 5 ans. 15% de croissance organique.",
        "Le saut ann\u00e9e 3 vient d\u2019une acquisition mid-market. La marge remonte gr\u00e2ce \u00e0 l\u2019ARR plateforme.")
td = [
    ["", "Constitution", "Ann\u00e9e 1\u20132", "Ann\u00e9e 3", "Ann\u00e9e 5"],
    ["CA groupe", "17M", "22M", "55M", "80M"],
    ["dont ARR", "~0", "2M", "7M", "20M"],
    ["EBITDA", "2,5M", "3,5M", "7M", "16M"],
    ["Marge EBITDA", "15%", "16%", "13%", "20%"],
    ["\u00c9v\u00e9nement cl\u00e9", "4 acquisitions fondatrices", "Entr\u00e9e fonds PE + brique RH", "Acquisition mid-market (~30M CA)", "Croissance organique + ARR"],
]
rows, cols = len(td), len(td[0])
ts = s.shapes.add_table(rows, cols, M, y, CW, Inches(2.5)); tbl = ts.table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c); cell.text = td[r][c]
        for p in cell.text_frame.paragraphs:
            for rn in p.runs:
                rn.font.size = Pt(11); rn.font.color.rgb = INK
                if r == 0: rn.font.bold = True; rn.font.color.rgb = INK_MUTED
                if c == 0 and r > 0: rn.font.bold = True
                if r == 2: rn.font.bold = True; rn.font.color.rgb = BRONZE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HL_BG if r == 2 else (PAPER_WARM if r == 0 else WHITE)
punch(s, y+Inches(2.7), "La marge EBITDA baisse temporairement \u00e0 l\u2019ann\u00e9e 3 (co\u00fbts d\u2019int\u00e9gration) puis converge vers 20% gr\u00e2ce \u00e0 la mont\u00e9e en charge de l\u2019ARR \u2014 la ligne qui fait la valorisation.")

# ===== S12 L'\u00c9QUIPE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte V", "L\u2019\u00e9quipe", "Une \u00e9quipe qui a fait le chemin \u2014 pas une qui le th\u00e9orise.",
        "Construction de groupe \u00b7 Transformation industrielle \u00b7 R&D IA en production.", True)
cw4 = (CW - Inches(0.6)) / 4
for i, (nm, rl, bio) in enumerate([
    ("Yvan Chabanne", "CEO du groupe", "Ex-DG Scalian (ESN ~5 000 personnes). Construction et pilotage d\u2019un groupe multi-filiales dans le conseil tech et l\u2019ing\u00e9nierie."),
    ("William Le Ferrand", "CTO groupe \u00b7 CEO Hypsis", "Ex-CTO Sogeclair (ETI a\u00e9ronautique, 1 300 pers.). Transformation digitale d\u2019une ETI industrielle cot\u00e9e. Cr\u00e9ateur de la plateforme Hypsis."),
    ("Emmanuel Benazera", "CSO groupe \u00b7 CEO Jolibrain", "10 ans de R&D IA en production. 5 PhDs, publications NeurIPS / ICML / AAAI. Clients : Airbus, Dassault, SNCF, CNES. 50+ GPUs."),
    ("Florian Lepage", "CRO groupe", "Ex-COO d\u2019un \u00e9diteur de logiciel RH. Exp\u00e9rience du scale-up commercial et de la structuration d\u2019une offre r\u00e9currente."),
]):
    left = M + i*(cw4+Inches(0.2))
    cbox(s, left, y, cw4, Inches(2.2), CARD_DARK, GOLD)
    t = txb(s, left+Inches(0.12), y+Inches(0.1), cw4-Inches(0.24), Inches(2.0)); tf = t.text_frame; tf.word_wrap = True
    sr(tf, nm, 12, True, c=PAPER); ap(tf, rl, 10, True, c=GOLD); ap(tf, bio, 9, c=TXT_L)
y += Inches(2.5)
for i, ph in enumerate(["CEO conseil strat\u00e9gique Data & IA", "CEO verticale industrielle", "CEO transformation humaine", "\u2026"]):
    left = M + i*(cw4+Inches(0.2))
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, cw4, Inches(0.8))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD_DARK; sh.line.color.rgb = INK_MUTED; sh.line.width = Pt(1); sh.line.dash_style = 4
    t = txb(s, left, y+Inches(0.15), cw4, Inches(0.5)); sr(t.text_frame, ph, 10, c=INK_MUTED, al=PP_ALIGN.CENTER)
punch(s, y+Inches(1.0), "L\u2019\u00e9quipe fondatrice couvre la construction de groupe, la transformation industrielle, la R&D IA en production et le scale-up commercial. Chaque fondateur qui rejoint le groupe compl\u00e8te le puzzle.", True)

# ===== S13 POSITIONNEMENT =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte V", "Positionnement", "Ce que chaque acteur fait. Ce qui lui manque.",
        "Nous assemblons les meilleures briques existantes plut\u00f4t que de tout construire from scratch.")
pd = [
    ["Acteur", "Ce qu\u2019il fait", "Ce qui lui manque"],
    ["Accenture / Capgemini", "Transformation IA grands groupes", "Trop cher pour ETI, pas souverain, lock-in"],
    ["ESN mid-market", "Int\u00e9gration tech, data", "Pas de vision strat\u00e9gique, pas de mod\u00e8les IA propri\u00e9taires"],
    ["Cabinets conseil (McKinsey, BCG)", "Strat\u00e9gie, diagnostic", "Pas de d\u00e9ploiement, pas d\u2019infra, partent apr\u00e8s les slides"],
    ["\u00c9diteurs IA (Dataiku, Palantir)", "Plateforme data/IA", "Lock-in, cloud, pas d\u2019accompagnement humain"],
    ["Cabinets conseil Data & IA", "Strat\u00e9gie IA, Master Plan, CDO Office", "Pas de plateforme, pas de mod\u00e8les IA, d\u00e9pendants d\u2019int\u00e9grateurs"],
    ["Le groupe", "Strat\u00e9gie + Data + IA + Infra souveraine + Humain", "Int\u00e9gr\u00e9, souverain, pas de lock-in, propri\u00e9t\u00e9 client"],
]
rows, cols = len(pd), len(pd[0])
ts = s.shapes.add_table(rows, cols, M, y, CW, Inches(3.0)); tbl = ts.table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c); cell.text = pd[r][c]
        for p in cell.text_frame.paragraphs:
            for rn in p.runs:
                rn.font.size = Pt(11); rn.font.color.rgb = INK
                if r == 0: rn.font.bold = True; rn.font.color.rgb = INK_MUTED
                if r == rows-1: rn.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = HL_BG if r == rows-1 else (PAPER_WARM if r == 0 else WHITE)

# ===== S14 POURQUOI NOUS REJOINDRE =====
s = prs.slides.add_slide(BL)
y = hdr(s, "Acte VI", "Pourquoi nous rejoindre", "Vous gardez votre entreprise. Vous acc\u00e9dez \u00e0 un groupe.",
        "Ce que le groupe change pour un fondateur qui nous rejoint.", True)
cw3 = (CW - Inches(0.4)) / 3
for i, (t_, b_) in enumerate([
    ("Votre identit\u00e9", "Votre CEO, votre P&L, votre \u00e9quipe, vos clients. Le groupe ne r\u00e9\u00e9crit pas ce qui fonctionne."),
    ("Votre front commercial", "Programmes de transformation \u00e0 100\u2013500K vendus au niveau COMEX. Vos comp\u00e9tences int\u00e9gr\u00e9es dans une offre plus large, plus visible, vendue plus cher."),
    ("Du conseil \u00e0 l\u2019impact", "Vous ne sous-traitez plus \u00e0 des ESN. Vous gardez le contr\u00f4le de la strat\u00e9gie ET vous d\u00e9ployez avec la plateforme, les mod\u00e8les, l\u2019infrastructure du groupe."),
]):
    bc = GOLD if i == 2 else None
    ctxt(s, M+i*(cw3+Inches(0.2)), y, cw3, Inches(1.6), t_, b_, bg_c=CARD_DARK, bc=bc, tc=PAPER, bc2=TXT_L)
y += Inches(1.9); cw2 = (CW - Inches(0.3)) / 2
for i, (t_, b_) in enumerate([
    ("Participation \u00e0 la valeur", "BSPCE au niveau de la holding. Roll-over en equity. La croissance d\u2019une BU b\u00e9n\u00e9ficie \u00e0 tous. Les int\u00e9r\u00eats sont align\u00e9s."),
    ("Protection, pas dilution", "P&L s\u00e9par\u00e9s = chaque entit\u00e9 valorisable ind\u00e9pendamment. Le groupe peut se vendre \u00ab \u00e0 l\u2019appartement \u00bb ou en bloc. Les deux chemins restent ouverts."),
]):
    ctxt(s, M+i*(cw2+Inches(0.3)), y, cw2, Inches(1.3), t_, b_, bg_c=CARD_DARK, tc=PAPER, bc2=TXT_L)

# ===== S15 CONTACT =====
s = prs.slides.add_slide(BL); bg(s, INK); y = Inches(0.8)
t = txb(s, M, y, CW, Inches(0.3)); sr(t.text_frame, "Acte VI", 9, True, c=GOLD, al=PP_ALIGN.CENTER); y += Inches(0.3)
t = txb(s, M, y, CW, Inches(0.3)); sr(t.text_frame, "PROCHAINE \u00c9TAPE", 11, True, c=GOLD, al=PP_ALIGN.CENTER); y += Inches(0.4)
t = txb(s, Inches(2), y, Inches(9), Inches(0.7)); sr(t.text_frame, "Construisons ensemble.", 32, True, c=PAPER, al=PP_ALIGN.CENTER); y += Inches(0.8)
t = txb(s, Inches(2.5), y, Inches(8), Inches(0.6))
sr(t.text_frame, "Le groupe se construit maintenant. Les briques techniques existent. L\u2019\u00e9quipe est en place. La pi\u00e8ce manquante, c\u2019est peut-\u00eatre vous.", 14, c=TXT_L, al=PP_ALIGN.CENTER)
y += Inches(0.9); cw3 = (CW - Inches(0.4)) / 3
for i, (t_, b_) in enumerate([
    ("Ce que nous avons", "\u2022 Plateforme souveraine en production\n\u2022 10 ans de R&D IA, mod\u00e8les d\u00e9ploy\u00e9s\n\u2022 30 consultants data, 6M CA\n\u2022 CEO ex-DG d\u2019un groupe de 5 000 pers."),
    ("Ce que nous cherchons", "\u2022 La brique conseil strat\u00e9gique Data & IA\n\u2022 Un acc\u00e8s direct aux COMEX\n\u2022 Une m\u00e9thodologie de transformation \u00e9prouv\u00e9e\n\u2022 Un fondateur qui veut aller plus loin"),
    ("Prochaine \u00e9tape", "\u2022 \u00c9change informel entre fondateurs\n\u2022 D\u00e9monstration de la plateforme\n\u2022 Exploration des synergies"),
]):
    bc = GOLD if i == 2 else None
    ctxt(s, M+i*(cw3+Inches(0.2)), y, cw3, Inches(2.0), t_, b_, bg_c=CARD_DARK, bc=bc, tc=PAPER, bc2=TXT_L)
y += Inches(2.3)
t = txb(s, Inches(4), y, Inches(5), Inches(0.4)); sr(t.text_frame, "ANTERIQ", 16, True, c=GOLD, al=PP_ALIGN.CENTER)
t = txb(s, Inches(4), y+Inches(0.4), Inches(5), Inches(0.3)); sr(t.text_frame, "Confidentiel \u2014 Avril 2026", 10, c=INK_MUTED, al=PP_ALIGN.CENTER)

out = "/home/wleferrand/dev/hypsis/hypsis/landing/group/anteriq-pitch.pptx"
prs.save(out)
print(f"Saved to {out}")
